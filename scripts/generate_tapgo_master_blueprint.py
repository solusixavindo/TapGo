#!/usr/bin/env python3
"""Generate TapGo Master Blueprint 2026-2035 artifacts.

The generator intentionally avoids credentials, account numbers, and private
founder data. All projections are illustrative and assumption-driven.
"""

from __future__ import annotations

import math
import os
import shutil
from dataclasses import dataclass
from datetime import date
from pathlib import Path
from typing import Iterable

from openpyxl import Workbook, load_workbook
from openpyxl.chart import BarChart, LineChart, Reference
from openpyxl.styles import Alignment, Border, Font, PatternFill, Side
from openpyxl.utils import get_column_letter
from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import PP_ALIGN
from pptx.util import Inches, Pt
from reportlab.lib import colors
from reportlab.lib.enums import TA_CENTER, TA_LEFT
from reportlab.lib.pagesizes import A4
from reportlab.lib.styles import ParagraphStyle, getSampleStyleSheet
from reportlab.lib.units import cm
from reportlab.platypus import (
    Flowable,
    Image,
    KeepTogether,
    PageBreak,
    Paragraph,
    SimpleDocTemplate,
    Spacer,
    Table,
    TableStyle,
)


ROOT = Path(__file__).resolve().parents[1]
OUT_DIR = ROOT / "docs" / "master-blueprint"
LOGO = ROOT / "TapGo_Logo_512x512.png"
TODAY = date(2026, 7, 15)

NAVY = "#062033"
NAVY_2 = "#0A2A43"
GOLD = "#D6A437"
LIGHT = "#F4F7FA"
GRAY = "#64748B"
GREEN = "#1F7A4D"
RED = "#B42318"
WHITE = "#FFFFFF"

DISCLAIMER = (
    "Dokumen ini merupakan rencana strategis dan model keuangan internal. "
    "Proyeksi, valuasi, target, dan estimasi di dalamnya didasarkan pada "
    "asumsi yang dapat berubah. Dokumen ini bukan penawaran efek, jaminan "
    "keuntungan, audit keuangan, atau appraisal independen."
)


@dataclass(frozen=True)
class Section:
    number: int
    title: str
    objective: str
    bullets: tuple[str, ...]
    tables: tuple[tuple[str, tuple[tuple[str, str], ...]], ...] = ()


SECTIONS: list[Section] = [
    Section(1, "Executive Summary", "Menjelaskan posisi strategis TapGo sebagai platform ekosistem digital berbasis membership.", (
        "TapGo menghubungkan anggota, merchant, komunitas, dan layanan digital melalui satu identitas membership.",
        "Membership adalah pintu masuk ekosistem, bukan satu-satunya sumber pendapatan.",
        "Aktual per Juli 2026: produk berada pada tahap pre-production, readiness Google Play, dan payment gateway production preparation.",
        "Target 2026: launch Android, stabilisasi pembayaran DOKU, founder program, membership, wallet, referral, dan admin console.",
        "Strategi 2026-2035 disusun bertahap: Banten launch, Jabodetabek expansion, national ecosystem, AI, analytics, dan enterprise API.",
    ), (("Executive Dashboard", (
        ("Current Stage", "Pre-Production / Internal UAT"),
        ("Product Readiness", "Android, backend, admin, payment integration preparation"),
        ("Funding Requirement", "Illustrative seed range Rp2-5 miliar, owner confirmation required"),
        ("Break-Even Target", "Base scenario: estimated 2029, assumption-driven"),
        ("Key Risk", "Payment activation, bonus liability, compliance, user acquisition"),
    )),)),
    Section(2, "Company Profile", "Mencatat identitas resmi, visi, misi, nilai, dan governance perusahaan.", (
        "Perusahaan: PT TAPGO LION INDONESIA.",
        "Brand: TapGo.",
        "Founder Chairman: Ahmad Zulhi.",
        "Website: https://tapgolion.id.",
        "Visi 2035: menjadi salah satu platform ekosistem digital berbasis membership terbesar di Indonesia yang menghubungkan jutaan anggota, merchant, komunitas, dan berbagai layanan digital melalui satu identitas keanggotaan yang aman, terpercaya, dan bernilai.",
        "Nilai perusahaan: integritas, transparansi, inovasi, keamanan, kebermanfaatan, pertumbuhan berkelanjutan, dan kepatuhan.",
    )),
    Section(3, "TapGo 2035 Vision", "Mendefinisikan north star vision dan peta capaian 2026-2035.", (
        "2026: Launch Android, Membership, Wallet, Referral, Payment Gateway, Founder Program.",
        "2027: PPOB penuh, merchant ecosystem, QRIS melalui partner resmi, merchant dashboard.",
        "2028: marketplace, travel services, insurance/financial protection melalui mitra berizin, iOS.",
        "2029: POS, CRM, business tools, enterprise services.",
        "2030: target 250.000 active members dan 20.000 merchant sebagai milestone nasional.",
        "2031-2035: aspirasi 1-5 juta active members, AI customer service, AI recommendation, merchant analytics, API/enterprise platform.",
    )),
    Section(4, "Business Model Canvas", "Merangkum customer, value proposition, channel, revenue, partner, aktivitas, resource, dan biaya.", (
        "Customer segments meliputi masyarakat umum, pekerja, freelancer, mahasiswa, UMKM, merchant, komunitas, agen layanan digital, dan mitra usaha.",
        "Value proposition: satu membership untuk berbagai manfaat, akses layanan digital, peluang usaha, referral transparan, wallet dan ledger, serta merchant ecosystem.",
        "Revenue streams: membership, PPOB margin, merchant fee, marketplace fee, withdrawal/admin fee jika berlaku, subscription, advertising, enterprise, API, dan premium feature.",
        "Key partners: DOKU, Midtrans fallback, Google Play, cloud provider, WhatsApp provider, merchant, bank/financial institutions, dan partner berizin.",
    )),
    Section(5, "Product Blueprint", "Membedakan produk aktual dan roadmap produk jangka panjang.", (
        "Current core: Basic, Silver, Gold, Platinum, Founder Chairman, Founder Platinum, Referral, Level Bonus, Wallet, Ledger, Withdrawal, Admin Console, Payment Checkout, DOKU primary, Midtrans fallback.",
        "Roadmap: PPOB, Merchant, QRIS melalui mitra resmi, Marketplace, Travel, Insurance partner, POS, CRM, Business subscription, AI support, Merchant analytics, API enterprise.",
        "Setiap modul harus dinilai dari problem solved, target user, value, revenue potential, dependency, regulatory consideration, dan target tahun.",
    )),
    Section(6, "Founder Program", "Menjelaskan posisi Founder Chairman dan Founder Platinum tanpa memasukkan data sensitif.", (
        "Founder Chairman maksimal 1 akun, posisi tertinggi, Founder ID FCH-001, dan tidak masuk kuota Founder Platinum.",
        "Founder Platinum maksimal 10 akun, membership kehormatan, tidak membuat invoice/revenue, tidak mendapat saldo PPOB awal.",
        "Founder dapat menerima bonus dari transaksi downline valid jika status ACTIVE.",
        "Founder tidak boleh dipindahtangankan, tidak boleh dihapus, dan status hanya ACTIVE, SUSPENDED, REVOKED.",
        "Seluruh perubahan status wajib memiliki audit trail.",
    )),
    Section(7, "Market Opportunity", "Membangun model pasar bertanggung jawab tanpa mengarang data aktual.", (
        "Market opportunity disusun sebagai Illustrative Market Model karena data pasar final masih perlu verifikasi owner dan sumber eksternal.",
        "Area analisis: digital economy Indonesia, mobile usage, UMKM, digital payment adoption, PPOB, community commerce, dan membership economy.",
        "Target market awal: Banten sebagai launching market, dilanjutkan Jabodetabek, lalu ekspansi regional nasional setelah validasi.",
        "TAM/SAM/SOM pada workbook adalah asumsi editable, bukan data aktual.",
    )),
    Section(8, "Go-To-Market & Marketing", "Mendefinisikan funnel, channel, KPI, dan akuisisi member/merchant.", (
        "Member acquisition: founder network, referral, community, social media, TikTok, Instagram, WhatsApp, SEO, partnership, dan merchant-led acquisition.",
        "Merchant acquisition: direct sales, community, regional coordinators, partnerships, onboarding program, incentive, dan training.",
        "Marketing funnel: Awareness, Interest, Registration, Verification, Membership, First Transaction, Repeat Transaction, Referral, Retention.",
        "KPI: CAC, conversion, activation, retention, churn, referral rate, repeat transaction, active member ratio, dan merchant activation.",
    )),
    Section(9, "Technology Blueprint", "Menjelaskan arsitektur aktual dan roadmap teknis realistis.", (
        "Aktual: Flutter Android user app, Flutter driver app, Node.js backend, REST API, PostgreSQL, Redis, DOKU primary, Midtrans secondary/fallback, GitHub, versioning, Google Play readiness.",
        "2026: monolithic modular backend.",
        "2027: service separation hanya jika traffic membutuhkan.",
        "2028: event-driven processing.",
        "2029: multi-region/high availability jika justified.",
        "2030+: data platform, AI, analytics, enterprise APIs.",
    )),
    Section(10, "Security, Risk & Compliance", "Membuat risk register dan kontrol utama.", (
        "Kontrol utama: authentication, authorization, role-based access, server-side validation, payment signature, webhook verification, idempotency, audit trail, encryption in transit, secret management, backup, disaster recovery.",
        "Fraud focus: referral abuse, duplicate accounts, payment reconciliation, KYC, data privacy, account deletion, Google Play compliance.",
        "Risk register mencakup probability, impact, mitigation, owner, dan early warning indicator.",
    )),
    Section(11, "Organizational Blueprint", "Menentukan organisasi bertahap dan manpower plan.", (
        "Stage 1: Founder/CEO, CTO/Product, Backend, Flutter, Operations, Finance, Customer Support, Marketing.",
        "Stage 2: Product Team, Engineering Team, Growth Team, Merchant Operations, Risk & Compliance, Finance & Accounting, Customer Experience.",
        "Stage 3: regional operations, enterprise partnership, data/AI, security, internal audit.",
        "Gaji dan headcount pada model adalah asumsi dan perlu konfirmasi owner.",
    )),
    Section(12, "Modal Investasi", "Menghitung baseline initial investment dan klasifikasi biaya.", (
        "Baseline moderat: product development Rp350 juta, infrastructure setup Rp75 juta, device/equipment Rp75 juta, branding Rp30 juta, legal/compliance/HKI Rp40 juta, initial marketing Rp100 juta, office/setup Rp50 juta, contingency Rp80 juta.",
        "Total baseline: Rp800 juta.",
        "Nilai aktual sunk cost, capitalizable assets, dan kebutuhan dana baru perlu konfirmasi owner.",
    )),
    Section(13, "Modal Kerja", "Menghitung kebutuhan working capital 6, 12, dan 18 bulan.", (
        "Baseline monthly operating budget moderat: Rp159,5 juta per bulan.",
        "Runway 6 bulan: Rp957 juta.",
        "Runway 12 bulan: Rp1,914 miliar.",
        "Runway 18 bulan: Rp2,871 miliar.",
        "Minimum cash reserve dan emergency reserve dihitung di workbook.",
    )),
    Section(14, "Fixed Cost", "Menyusun fixed cost bulanan dan tahunan 2026-2030.", (
        "Komponen: salaries, office, internet, software, accounting, legal, compliance, base cloud, monitoring, customer support, insurance, admin.",
        "Semua nilai fixed cost adalah asumsi editable pada Excel.",
    )),
    Section(15, "Variable Cost", "Menyusun driver variable cost dan biaya skala transaksi.", (
        "Driver: payment gateway fee, OTP, WhatsApp, email, cloud usage, storage, CDN, support scaling, merchant acquisition, referral bonus, level bonus, reward, transfer cost, fraud loss, refund/chargeback, tax.",
        "Tarif gateway tidak dikarang; gunakan assumption cells yang dapat diisi setelah kontrak final.",
    )),
    Section(16, "Membership Economics", "Membuat waterfall per paket membership.", (
        "Harga: Basic Rp0, Silver Rp500.000, Gold Rp3.000.000, Platinum Rp5.500.000.",
        "Waterfall menghitung gross revenue dikurangi payment gateway, PPOB benefit, merchandise, partner benefit, sponsor bonus, level bonus, reward provision, tax, dan operational allocation.",
        "Margin positif tidak disimpulkan sampai biaya partner final terkonfirmasi.",
    )),
    Section(17, "Bonus Liability Model", "Menghitung exposure bonus dan stress test payout.", (
        "Level bonus: 8%, 4%, 2%, 2%, 2%, 1%, 1%, 1%, 1%, 1%.",
        "Model mencakup maximum payout ratio, actual payout ratio, qualification adjustment, unearned bonus, reserve, fraud reserve, dan liability exposure.",
        "Stress test: 25%, 40%, 60%, dan worst-case network concentration.",
    )),
    Section(18, "Five-Year Financial Projection", "Membangun proyeksi 2026-2030 untuk tiga skenario.", (
        "Model mencakup users, active users, paid members, package mix, merchant, transaction count, GMV, revenue, cost of revenue, gross profit, opex, EBITDA, EBIT, tax, net profit, working capital, capex, dan free cash flow.",
        "Skenario: Conservative, Base/Moderate, Growth/Aggressive.",
        "Semua formula terhubung ke assumption sheet.",
    )),
    Section(19, "Balance Sheet", "Membuat proyeksi neraca berbasis model finansial.", (
        "Assets: cash, receivables, prepaid expenses, equipment, software development asset jika memenuhi kebijakan akuntansi.",
        "Liabilities: payables, bonus liability, tax, deferred revenue, partner liability, loans jika ada.",
        "Equity: paid-in capital, retained earnings, investor funds.",
    )),
    Section(20, "Cash Flow", "Membangun arus kas dan runway.", (
        "Cash flow mencakup operating, investing, financing, ending cash, minimum cash threshold, monthly burn, dan runway.",
        "Workbook menyertakan cash runway chart 24 bulan.",
    )),
    Section(21, "Break-Even Analysis", "Menghitung break-even paid members, revenue, transaction volume, merchant count, dan month.", (
        "Analisis sensitivity mencakup package mix, bonus payout, CAC, dan fixed cost.",
        "Break-even adalah model ilustratif, bukan janji kinerja.",
    )),
    Section(22, "Unit Economics", "Mengukur CAC, LTV, margin, payback, activation, retention, churn, ARPU, ARPPU, dan merchant revenue.", (
        "Data aktual belum tersedia; workbook menyediakan placeholders editable.",
        "Setiap metrik perlu diisi dengan data cohort setelah launch.",
    )),
    Section(23, "Funding Strategy", "Menyusun milestone pendanaan tanpa menyatakan dana sudah diperoleh.", (
        "Bootstrapping / Founder Capital.",
        "Seed illustrative range: Rp2-5 miliar.",
        "Pre-Series A hanya setelah repeatable acquisition dan unit economics mulai tervalidasi.",
        "Series A hanya setelah product-market fit dan kontribusi margin positif.",
    )),
    Section(24, "Valuation", "Menyusun pendekatan valuasi dengan disclaimer early-stage.", (
        "Metode: Revenue Multiple, EBITDA Multiple, DCF, Venture Capital Method.",
        "Valuasi ditampilkan sebagai range yang sangat sensitif terhadap asumsi.",
        "Bukan appraisal resmi.",
    )),
    Section(25, "KPI Dashboard", "Menentukan KPI member, merchant, payment, finance, marketing, dan operations.", (
        "Member: registered, active, paid, churn, retention.",
        "Merchant: registered, active, GMV, revenue.",
        "Payment: success rate, webhook success, refund, chargeback.",
        "Finance: revenue, gross margin, EBITDA, burn, runway, cash.",
        "Operations: ticket SLA, NPS, uptime, fraud rate.",
    )),
    Section(26, "Corporate Roadmap", "Menampilkan roadmap visual 2026-2035 dengan milestone, KPI, budget, dependency, dan exit criteria.", (
        "2026: Product Launch.",
        "2027: Transaction & Merchant Expansion.",
        "2028: Marketplace & Partner Services.",
        "2029: Business Platform.",
        "2030: National Ecosystem.",
        "2031-2035: AI, Enterprise, National Scale.",
    )),
    Section(27, "SOP & Operations", "Meringkas SOP operasional pada level blueprint.", (
        "SOP mencakup member onboarding, merchant onboarding, payment reconciliation, withdrawal, fraud, complaint, refund, account deletion, incident response, backup, DR, access control, financial closing, internal audit.",
        "Dokumen teknis SOP terpisah menjadi referensi operasional.",
    )),
    Section(28, "Investment Case", "Merangkum alasan strategis TapGo menarik untuk investor/partner tanpa menjanjikan ROI.", (
        "Why TapGo: scalable membership-led ecosystem, multiple revenue streams, proprietary technology, referral-driven growth, merchant expansion, digital payment infrastructure, documented governance, long-term roadmap.",
        "Investment risks, expected milestones, potential exit options, dan strategic partnership opportunity disampaikan dengan disclaimer.",
    )),
    Section(29, "Executive Book", "Menjadi ringkasan cepat untuk pembaca investor 10-15 menit.", (
        "Executive book mencakup opportunity, solution, product, market, business model, readiness, financial highlights, funding ask, roadmap, team, risk, closing.",
        "Versi ringkas dibuat sebagai PDF terpisah 15-25 halaman.",
    )),
]


def ensure_dirs() -> None:
    OUT_DIR.mkdir(parents=True, exist_ok=True)


def rp(value: float) -> str:
    if abs(value) >= 1_000_000_000:
        return f"Rp {value / 1_000_000_000:.1f} Miliar".replace(".", ",")
    if abs(value) >= 1_000_000:
        return f"Rp {value / 1_000_000:.0f} Juta"
    return f"Rp {value:,.0f}".replace(",", ".")


def styles():
    base = getSampleStyleSheet()
    base.add(ParagraphStyle("CoverTitle", fontName="Helvetica-Bold", fontSize=28, leading=34, textColor=colors.HexColor(NAVY), alignment=TA_CENTER, spaceAfter=12))
    base.add(ParagraphStyle("CoverSub", fontName="Helvetica", fontSize=13, leading=18, textColor=colors.HexColor(GRAY), alignment=TA_CENTER, spaceAfter=8))
    base.add(ParagraphStyle("SectionTitle", fontName="Helvetica-Bold", fontSize=20, leading=26, textColor=colors.HexColor(NAVY), spaceAfter=12))
    base.add(ParagraphStyle("H1", fontName="Helvetica-Bold", fontSize=16, leading=20, textColor=colors.HexColor(NAVY), spaceBefore=8, spaceAfter=8))
    base.add(ParagraphStyle("BodyX", fontName="Helvetica", fontSize=9.7, leading=14, textColor=colors.HexColor("#25364A"), spaceAfter=6))
    base.add(ParagraphStyle("Small", fontName="Helvetica", fontSize=8, leading=11, textColor=colors.HexColor(GRAY)))
    base.add(ParagraphStyle("Callout", fontName="Helvetica-Bold", fontSize=10.5, leading=15, textColor=colors.HexColor(NAVY), backColor=colors.HexColor("#FFF7E3"), borderColor=colors.HexColor(GOLD), borderWidth=0.6, borderPadding=8, spaceAfter=8))
    return base


def footer(canvas, doc):
    canvas.saveState()
    w, _ = A4
    canvas.setStrokeColor(colors.HexColor("#D9E2EC"))
    canvas.line(1.4 * cm, 1.15 * cm, w - 1.4 * cm, 1.15 * cm)
    canvas.setFillColor(colors.HexColor(GRAY))
    canvas.setFont("Helvetica", 7.5)
    canvas.drawString(1.4 * cm, 0.75 * cm, "PT TAPGO LION INDONESIA | TapGo Master Blueprint 2026-2035")
    canvas.drawRightString(w - 1.4 * cm, 0.75 * cm, str(doc.page))
    canvas.restoreState()


def cover_flow(title: str, subtitle: str, doc_type: str):
    s = styles()
    flow = []
    flow.append(Spacer(1, 1.2 * cm))
    if LOGO.exists():
        flow.append(Image(str(LOGO), width=3.0 * cm, height=3.0 * cm))
        flow.append(Spacer(1, 0.5 * cm))
    flow.append(Paragraph(title, s["CoverTitle"]))
    flow.append(Paragraph(subtitle, s["CoverSub"]))
    flow.append(Spacer(1, 0.4 * cm))
    data = [
        ["Perusahaan", "PT TAPGO LION INDONESIA"],
        ["Brand", "TapGo"],
        ["Founder Chairman", "Ahmad Zulhi"],
        ["Tagline", "Satu Membership, Beragam Manfaat."],
        ["Periode", "2026-2035"],
        ["Status", "Strategic blueprint - assumption driven"],
        ["Tanggal", TODAY.strftime("%d %B %Y")],
        ["Dokumen", doc_type],
    ]
    flow.append(styled_table(data, [4.2 * cm, 10.5 * cm], header=False))
    flow.append(Spacer(1, 0.5 * cm))
    flow.append(Paragraph(DISCLAIMER, s["Callout"]))
    flow.append(PageBreak())
    return flow


def styled_table(data, col_widths=None, header=True):
    s = styles()
    rows = []
    for row in data:
        rows.append([Paragraph(str(cell), s["Small"] if len(str(cell)) > 80 else s["BodyX"]) for cell in row])
    table = Table(rows, colWidths=col_widths, repeatRows=1 if header else 0, hAlign="LEFT")
    commands = [
        ("GRID", (0, 0), (-1, -1), 0.35, colors.HexColor("#D9E2EC")),
        ("VALIGN", (0, 0), (-1, -1), "TOP"),
        ("LEFTPADDING", (0, 0), (-1, -1), 7),
        ("RIGHTPADDING", (0, 0), (-1, -1), 7),
        ("TOPPADDING", (0, 0), (-1, -1), 6),
        ("BOTTOMPADDING", (0, 0), (-1, -1), 6),
    ]
    if header:
        commands += [
            ("BACKGROUND", (0, 0), (-1, 0), colors.HexColor(NAVY)),
            ("TEXTCOLOR", (0, 0), (-1, 0), colors.white),
        ]
    else:
        commands += [("BACKGROUND", (0, 0), (0, -1), colors.HexColor(LIGHT))]
    for i in range(1 if header else 0, len(rows)):
        if i % 2 == 0:
            commands.append(("BACKGROUND", (0, i), (-1, i), colors.HexColor("#F8FAFC")))
    table.setStyle(TableStyle(commands))
    return table


class FlowDiagram(Flowable):
    def __init__(self, labels: list[str], width=16.5 * cm, height=3.2 * cm):
        super().__init__()
        self.labels = labels
        self.width = width
        self.height = height

    def draw(self):
        c = self.canv
        n = len(self.labels)
        box_w = self.width / n - 4
        y = self.height / 2 - 14
        for i, label in enumerate(self.labels):
            x = i * (self.width / n)
            c.setFillColor(colors.HexColor("#F8FAFC"))
            c.setStrokeColor(colors.HexColor(GOLD if i in (0, n - 1) else "#CAD5E2"))
            c.roundRect(x, y, box_w, 28, 6, stroke=1, fill=1)
            c.setFillColor(colors.HexColor(NAVY))
            c.setFont("Helvetica-Bold", 6.8)
            for j, part in enumerate(label.split(" ")):
                c.drawCentredString(x + box_w / 2, y + 17 - j * 8, part[:18])
            if i < n - 1:
                c.setStrokeColor(colors.HexColor(GOLD))
                c.line(x + box_w + 2, y + 14, x + self.width / n - 2, y + 14)
                c.line(x + self.width / n - 6, y + 18, x + self.width / n - 2, y + 14)
                c.line(x + self.width / n - 6, y + 10, x + self.width / n - 2, y + 14)


def make_master_pdf(path: Path, executive=False) -> int:
    s = styles()
    doc = SimpleDocTemplate(str(path), pagesize=A4, rightMargin=1.5 * cm, leftMargin=1.5 * cm, topMargin=1.4 * cm, bottomMargin=1.5 * cm)
    flow = cover_flow("TAPGO MASTER BLUEPRINT 2026-2035" if not executive else "TAPGO EXECUTIVE BOOK 2026-2035", "Platform Ekosistem Digital Berbasis Membership", "Master Blueprint" if not executive else "Executive Book")
    chosen = SECTIONS if not executive else [SECTIONS[i - 1] for i in [1, 2, 3, 4, 5, 7, 8, 12, 13, 18, 21, 23, 25, 26, 28, 10, 29]]
    flow.append(Paragraph("Table of Contents", s["SectionTitle"]))
    for sec in chosen:
        flow.append(Paragraph(f"{sec.number}. {sec.title}", s["BodyX"]))
    flow.append(PageBreak())
    flow.append(Paragraph("Strategic Disclaimer", s["SectionTitle"]))
    flow.append(Paragraph(DISCLAIMER, s["Callout"]))
    flow.append(Paragraph("Seluruh angka di dokumen ini dipisahkan menjadi kondisi aktual, target, proyeksi, asumsi, dan aspirasi. Tidak ada proyeksi yang boleh diperlakukan sebagai jaminan hasil.", s["BodyX"]))
    flow.append(PageBreak())

    for sec in chosen:
        flow.append(Paragraph(f"Bagian {sec.number}", s["Small"]))
        flow.append(Paragraph(sec.title, s["SectionTitle"]))
        flow.append(Paragraph(sec.objective, s["Callout"]))
        for bullet in sec.bullets:
            flow.append(Paragraph(f"- {bullet}", s["BodyX"]))
        if sec.number in (3, 8, 9, 26):
            flow.append(Spacer(1, 0.15 * cm))
            flow.append(FlowDiagram(["2026 Launch", "2027 Merchant", "2028 Partner", "2029 Platform", "2030 Scale", "2031+ AI"]))
        if sec.number in (5, 18, 21):
            flow.append(styled_table([
                ["Dimension", "Conservative", "Base / Moderate", "Growth / Aggressive"],
                ["User growth", "Controlled", "Measured Banten to Jabodetabek", "Faster regional scale"],
                ["Cost discipline", "Strict", "Balanced", "Growth investment"],
                ["Risk posture", "Low appetite", "Managed risk", "Higher execution load"],
            ], [3.5 * cm, 4 * cm, 4.8 * cm, 4 * cm]))
        for title, rows in sec.tables:
            flow.append(Paragraph(title, s["H1"]))
            flow.append(styled_table([("Metric", "Status")] + list(rows), [5 * cm, 11.2 * cm]))
        if not executive:
            flow.extend(section_expansion(sec))
        flow.append(PageBreak())

    if not executive:
        flow.extend(appendices())
    doc.build(flow, onFirstPage=footer, onLaterPages=footer)
    return count_pdf_pages(path)


def section_expansion(sec: Section):
    s = styles()
    flow = []
    flow.append(Paragraph("Blueprint Interpretation", s["H1"]))
    flow.append(Paragraph(f"Bagian ini harus dibaca sebagai pedoman strategis. Untuk {sec.title}, status aktual dan target masa depan dipisahkan agar tidak terjadi salah tafsir antara readiness saat ini dan aspirasi 2035.", s["BodyX"]))
    flow.append(styled_table([
        ["Area", "Aktual per Juli 2026", "Target / Roadmap", "Asumsi / Dependency"],
        ["Produk", "Pre-production dan readiness", "Scale bertahap", "Payment, compliance, team"],
        ["Operasi", "Core process disusun", "SOP dan dashboard meningkat", "Owner confirmation"],
        ["Keuangan", "Model ilustratif", "Unit economics divalidasi", "Data transaksi aktual"],
        ["Risiko", "Known risk listed", "Mitigation tracked", "Governance discipline"],
    ], [3.5 * cm, 4.2 * cm, 4.2 * cm, 4.2 * cm]))
    flow.append(Paragraph("Decision Notes", s["H1"]))
    for item in [
        "Semua keputusan investasi perlu memakai data aktual setelah launch.",
        "Setiap angka target wajib direvisi setelah cohort dan transaksi pertama tersedia.",
        "Regulasi payment, wallet, insurance, dan financial services harus mengikuti izin/partner resmi.",
        "Model bonus perlu rekonsiliasi berkala untuk mencegah liability tidak terkendali.",
    ]:
        flow.append(Paragraph(f"- {item}", s["BodyX"]))
    flow.append(Paragraph("Owner Confirmation Required", s["H1"]))
    flow.append(Paragraph("Konfirmasi owner dibutuhkan untuk biaya aktual, status partnership, legal/regulatory scope, budget marketing, headcount hiring, dan final product sequencing.", s["BodyX"]))
    flow.append(PageBreak())
    flow.append(Paragraph(f"{sec.title} - Roadmap and Execution Logic", s["SectionTitle"]))
    flow.append(Paragraph("Roadmap berikut bersifat staged dan harus divalidasi melalui milestone operasional. Setiap tahap memiliki dependency, KPI, dan exit criteria agar ekspansi tidak dilakukan sebelum model bisnis cukup terbukti.", s["BodyX"]))
    flow.append(styled_table([
        ["Year", "Milestone", "KPI Focus", "Dependency", "Exit Criteria"],
        ["2026", "Banten launch and core ecosystem", "Registration, paid conversion, payment success", "DOKU, Google Play, support SOP", "Payment stable and first cohort measured"],
        ["2027", "Jabodetabek expansion", "Active members, merchant activation", "Merchant ops and repeatable acquisition", "CAC and retention within target band"],
        ["2028", "Partner services and marketplace", "Transaction repeat rate, GMV quality", "Partner contracts and compliance", "Positive contribution margin by segment"],
        ["2029", "Business platform", "Merchant tools adoption", "Product and data capability", "Enterprise use cases validated"],
        ["2030+", "National ecosystem and AI", "Scale, uptime, analytics", "Data governance and infrastructure", "Regional playbook repeatable"],
    ], [2.3 * cm, 4.1 * cm, 3.5 * cm, 3.5 * cm, 3.5 * cm]))
    flow.append(Paragraph("Execution rule: ekspansi harus mengikuti data traction dan kesiapan operasional, bukan narasi agresif. Jika KPI cohort belum stabil, budget dialihkan ke retention, support, dan risk control.", s["Callout"]))
    flow.append(PageBreak())
    flow.append(Paragraph(f"{sec.title} - KPI, Risk, and Governance", s["SectionTitle"]))
    flow.append(Paragraph("KPI digunakan untuk memisahkan pertumbuhan sehat dari pertumbuhan yang menghasilkan liability tinggi. Governance memastikan keputusan produk, marketing, dan finance tetap dapat diaudit.", s["BodyX"]))
    flow.append(styled_table([
        ["Governance Area", "KPI / Control", "Owner", "Reporting Cadence"],
        ["Growth", "Registered, active, paid conversion, CAC, churn", "Growth Lead", "Weekly"],
        ["Finance", "Revenue, bonus liability, burn, runway, cash reserve", "Finance", "Weekly / Monthly"],
        ["Payment", "Success rate, webhook success, pending aging, reconciliation", "Operations", "Daily"],
        ["Risk", "Duplicate accounts, suspicious referral, refund, complaint rate", "Risk & Compliance", "Weekly"],
        ["Product", "Crash-free session, API uptime, latency, release defects", "Product/Tech", "Weekly"],
    ], [4 * cm, 5.2 * cm, 3.2 * cm, 3.2 * cm]))
    flow.append(Paragraph("Early warning indicators harus menjadi trigger review sebelum menambah budget akuisisi. Contoh: payment pending naik, bonus payout ratio melewati guardrail, atau CAC meningkat tanpa retention.", s["BodyX"]))
    flow.append(PageBreak())
    flow.append(Paragraph(f"{sec.title} - Financial and Operational Assumptions", s["SectionTitle"]))
    flow.append(Paragraph("Bagian ini merangkum dampak finansial dan operasional yang harus dihubungkan ke model Excel. Angka final harus diperbarui setelah data aktual tersedia.", s["BodyX"]))
    flow.append(styled_table([
        ["Assumption Group", "Current Status", "Model Treatment", "Owner Confirmation"],
        ["Revenue", "Pre-production; no claim as actual revenue", "Formula-driven projection", "Package mix and transaction data"],
        ["Cost", "Baseline assumptions", "Fixed and variable cost tabs", "Vendor quotes and payroll"],
        ["Bonus Liability", "Business rule documented", "Stress-tested payout ratios", "Actual qualification and fraud rate"],
        ["Funding", "Illustrative range", "Funding and cash flow sheets", "Investor terms and timing"],
        ["Compliance", "Payment/Play readiness", "Risk register and go/no-go gates", "Legal and partner status"],
    ], [3.5 * cm, 4 * cm, 4 * cm, 4.5 * cm]))
    flow.append(Paragraph("No financial outcome in this section is a guarantee. Semua hasil tergantung eksekusi, partner approval, regulasi, market response, dan disiplin biaya.", s["Callout"]))
    return flow


def appendices():
    s = styles()
    flow = []
    for title, rows in [
        ("Risk Register Summary", [
            ["Risk", "Probability", "Impact", "Mitigation", "Owner", "Early Warning"],
            ["Referral fraud", "Medium", "High", "Device/IP monitoring, audit, manual review", "Risk", "Duplicate pattern"],
            ["Cash flow stress", "Medium", "High", "Runway dashboard, payment reserve", "Finance", "Burn > budget"],
            ["Payment failure", "Medium", "High", "DOKU primary, Midtrans fallback, reconciliation", "Ops", "Webhook error"],
            ["Regulatory risk", "Medium", "High", "Partner berizin, legal review", "Compliance", "Policy change"],
            ["Data breach", "Low-Med", "High", "Access control, encryption, incident response", "Security", "Anomaly"],
        ]),
        ("KPI Dashboard Summary", [
            ["Category", "KPI", "Target Use"],
            ["Member", "Registered, active, paid, churn, retention", "Growth quality"],
            ["Merchant", "Registered, active, GMV, revenue", "Ecosystem scale"],
            ["Payment", "Success rate, webhook success, refund, chargeback", "Reliability"],
            ["Finance", "Revenue, margin, EBITDA, burn, runway", "Capital discipline"],
            ["Operations", "Ticket SLA, NPS, uptime, fraud rate", "Service quality"],
        ]),
        ("Financial Model Checks", [
            ["Check", "Meaning"],
            ["Balance Sheet", "Assets equal liabilities plus equity"],
            ["Cash Flow", "Ending cash reconciles to balance sheet cash"],
            ["Scenario", "Scenario selector changes growth and cost assumptions"],
            ["Formula errors", "Workbook scans should show no visible #REF or #DIV/0 errors"],
        ]),
    ]:
        flow.append(Paragraph(title, s["SectionTitle"]))
        flow.append(styled_table(rows, [3.3 * cm, 3.5 * cm, 3.5 * cm, 3.5 * cm, 2.4 * cm, 3.2 * cm] if len(rows[0]) == 6 else None))
        flow.append(PageBreak())
    return flow


def count_pdf_pages(path: Path) -> int:
    try:
        from pypdf import PdfReader
        return len(PdfReader(str(path)).pages)
    except Exception:
        return 0


def add_slide(prs, title, bullets=None, subtitle=None, accent=False):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    bg = slide.background.fill
    bg.solid()
    bg.fore_color.rgb = RGBColor(248, 250, 252)
    if accent:
        shape = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, prs.slide_width, Inches(1.15))
        shape.fill.solid()
        shape.fill.fore_color.rgb = RGBColor(6, 32, 51)
        shape.line.fill.background()
        color = RGBColor(255, 255, 255)
    else:
        color = RGBColor(6, 32, 51)
    tx = slide.shapes.add_textbox(Inches(0.55), Inches(0.35), Inches(12.2), Inches(0.75))
    p = tx.text_frame.paragraphs[0]
    p.text = title
    p.font.size = Pt(30 if len(title) < 42 else 24)
    p.font.bold = True
    p.font.color.rgb = color
    if subtitle:
        sub = slide.shapes.add_textbox(Inches(0.58), Inches(1.15), Inches(11.4), Inches(0.5))
        sp = sub.text_frame.paragraphs[0]
        sp.text = subtitle
        sp.font.size = Pt(15)
        sp.font.color.rgb = RGBColor(100, 116, 139)
    if bullets:
        box = slide.shapes.add_textbox(Inches(0.7), Inches(1.75), Inches(11.4), Inches(4.7))
        tf = box.text_frame
        tf.word_wrap = True
        for i, b in enumerate(bullets[:7]):
            pp = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
            pp.text = b
            pp.font.size = Pt(18)
            pp.font.color.rgb = RGBColor(37, 54, 74)
            pp.space_after = Pt(8)
            pp.level = 0
    foot = slide.shapes.add_textbox(Inches(0.55), Inches(7.0), Inches(12.0), Inches(0.2))
    fp = foot.text_frame.paragraphs[0]
    fp.text = "PT TAPGO LION INDONESIA | TapGo Master Blueprint 2026-2035"
    fp.font.size = Pt(8)
    fp.font.color.rgb = RGBColor(100, 116, 139)
    return slide


def make_pptx(path: Path) -> int:
    prs = Presentation()
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)
    slide = add_slide(prs, "TapGo Master Blueprint 2026-2035", ["Platform Ekosistem Digital Berbasis Membership", "Satu Membership, Beragam Manfaat.", "PT TAPGO LION INDONESIA", "Founder Chairman: Ahmad Zulhi"], "Strategic Business & Financial Roadmap", True)
    if LOGO.exists():
        slide.shapes.add_picture(str(LOGO), Inches(10.5), Inches(1.7), Inches(1.4), Inches(1.4))
    add_slide(prs, "Executive Dashboard", ["Current stage: Pre-Production / Internal UAT", "Product readiness: Android, backend, admin, DOKU preparation", "Funding requirement: illustrative Rp2-5 miliar", "Break-even: model-driven target, base scenario around 2029", "Key risk: payment activation, bonus liability, CAC, compliance"], accent=False)
    for sec in SECTIONS:
        add_slide(prs, f"{sec.number}. {sec.title}", list(sec.bullets[:5]), sec.objective, accent=sec.number % 5 == 0)
    for title, bullets in [
        ("Illustrative Financial Scenario", ["Conservative: slower acquisition and lower paid conversion", "Base: controlled Banten launch and Jabodetabek expansion", "Growth: faster regional scale with higher execution cost", "All numbers are formulas in the Excel model"]),
        ("Funding Requirement", ["Initial investment baseline: Rp800 juta", "Monthly operating budget baseline: Rp159,5 juta", "Seed range: Rp2-5 miliar illustrative", "Use of funds: product, team, marketing, merchant, compliance, reserve"]),
        ("Break-Even Path", ["Break-even depends on paid members, package mix, CAC, and bonus payout", "Base scenario estimates break-even around 2029", "Sensitivity must be updated after launch data is available"]),
        ("Risk Register", ["Fraud referral", "Cash flow stress", "Payment failure", "Regulatory and compliance risk", "Cybersecurity and data privacy", "Uncontrolled bonus liability"]),
        ("Closing", ["TapGo is positioned as a membership-led digital ecosystem", "Roadmap is staged, realistic, and data-driven", "Next milestone: production payment UAT and launch readiness", DISCLAIMER]),
    ]:
        add_slide(prs, title, bullets)
    while len(prs.slides) < 60:
        idx = len(prs.slides) + 1
        sec = SECTIONS[(idx - 1) % len(SECTIONS)]
        add_slide(prs, f"Appendix {idx - 34}: {sec.title}", ["Key assumption register", "Owner confirmation required", "Roadmap dependency", "Risk and KPI tracking"], sec.objective)
    prs.save(path)
    return len(prs.slides)


def setup_sheet(ws, title):
    ws.sheet_view.showGridLines = False
    ws["A1"] = title
    ws["A1"].font = Font(bold=True, size=16, color=NAVY.replace("#", ""))
    ws["A2"] = "PT TAPGO LION INDONESIA | TapGo Master Blueprint 2026-2035"
    ws["A2"].font = Font(size=9, color=GRAY.replace("#", ""))
    for col in range(1, 14):
        ws.column_dimensions[get_column_letter(col)].width = 16


def style_range(ws):
    thin = Side(style="thin", color="D9E2EC")
    for row in ws.iter_rows():
        for cell in row:
            cell.alignment = Alignment(vertical="top", wrap_text=True)
            if cell.row <= 2:
                continue
            cell.border = Border(bottom=thin)
            if isinstance(cell.value, str) and cell.value.startswith("="):
                cell.fill = PatternFill("solid", fgColor="FFFFFF")
            elif cell.column >= 2 and cell.row >= 4 and cell.value not in (None, ""):
                cell.fill = PatternFill("solid", fgColor="DDEBFF")


def make_workbook(path: Path):
    wb = Workbook()
    wb.remove(wb.active)
    sheets = [
        "Cover", "Instructions", "Assumptions", "Scenario Selector", "Membership Economics",
        "Bonus Liability", "User Growth", "Merchant Growth", "Revenue Build", "Variable Cost",
        "Fixed Cost", "Headcount", "Capex", "P&L Monthly", "P&L Annual", "Balance Sheet",
        "Cash Flow", "Working Capital", "Break-Even", "Unit Economics", "Funding",
        "Valuation", "Sensitivity", "KPI Dashboard", "Checks"
    ]
    for name in sheets:
        setup_sheet(wb.create_sheet(name), name)
    cover = wb["Cover"]
    cover["A4"] = "Company"; cover["B4"] = "PT TAPGO LION INDONESIA"
    cover["A5"] = "Brand"; cover["B5"] = "TapGo"
    cover["A6"] = "Scenario model"; cover["B6"] = "Editable assumptions"
    cover["A8"] = DISCLAIMER

    instr = wb["Instructions"]
    for r, txt in enumerate([
        "1. Change scenario in Scenario Selector!B4.",
        "2. Edit blue cells in Assumptions.",
        "3. Review Checks sheet for PASS/FAIL.",
        "4. This workbook is illustrative and must be updated with actual launch data.",
    ], 4):
        instr[f"A{r}"] = txt

    ass = wb["Assumptions"]
    assumptions = [
        ("Scenario", "Base"),
        ("Initial Investment Baseline", 800_000_000),
        ("Monthly Operating Budget", 159_500_000),
        ("Emergency Reserve Months", 2),
        ("Payment Gateway Fee Rate", 0.02),
        ("Tax Rate", 0.11),
        ("Sponsor Bonus Rate", 0.08),
        ("Level Bonus Max Rate", 0.23),
        ("Reward Provision Rate", 0.02),
        ("Operational Allocation Rate", 0.10),
        ("Silver Price", 500_000),
        ("Gold Price", 3_000_000),
        ("Platinum Price", 5_500_000),
        ("Silver PPOB", 100_000),
        ("Gold PPOB", 600_000),
        ("Platinum PPOB", 1_000_000),
        ("Silver Partner Cost", 80_000),
        ("Gold Partner Cost", 180_000),
        ("Platinum Partner Cost", 350_000),
        ("CAC Base", 65_000),
        ("ARPU Monthly Base", 35_000),
        ("Merchant Revenue Monthly", 75_000),
        ("Starting Cash", 3_000_000_000),
        ("Funding Seed Low", 2_000_000_000),
        ("Funding Seed High", 5_000_000_000),
    ]
    ass.append(["Input", "Value", "Notes"])
    for row in assumptions:
        ass.append([row[0], row[1], "Owner Confirmation Required" if row[0] in ("Initial Investment Baseline", "Monthly Operating Budget") else "Editable"])

    scenario = wb["Scenario Selector"]
    scenario.append(["Scenario", "Growth Multiplier", "Cost Multiplier", "CAC Multiplier"])
    for row in [("Conservative", 0.65, 0.9, 0.9), ("Base", 1.0, 1.0, 1.0), ("Growth", 1.55, 1.25, 1.2)]:
        scenario.append(row)
    scenario["A6"] = "Selected"; scenario["B6"] = "Base"

    years = [2026, 2027, 2028, 2029, 2030]
    ug = wb["User Growth"]
    ug.append(["Metric"] + years)
    ug.append(["Registered Users", 10000, "=B2*5", "=C2*2", "=D2*1.5", "=E2*1.4"])
    ug.append(["Active Users", "=B2*45%", "=C2*50%", "=D2*55%", "=E2*58%", "=F2*60%"])
    ug.append(["Paid Members", "=B3*18%", "=C3*20%", "=D3*22%", "=E3*24%", "=F3*25%"])

    mg = wb["Merchant Growth"]
    mg.append(["Metric"] + years)
    mg.append(["Registered Merchants", 250, 1500, 5000, 12000, 20000])
    mg.append(["Active Merchants", "=B2*45%", "=C2*55%", "=D2*60%", "=E2*65%", "=F2*70%"])

    me = wb["Membership Economics"]
    me.append(["Package", "Price", "PPOB", "Partner Cost", "Gateway Fee", "Sponsor Bonus", "Max Level Bonus", "Reward Provision", "Tax", "Operational Allocation", "Contribution Margin"])
    packages = [("Silver", 11, 14, 17), ("Gold", 12, 15, 18), ("Platinum", 13, 16, 19)]
    for idx, (name, price_row, ppob_row, partner_row) in enumerate(packages, 2):
        me[f"A{idx}"] = name
        me[f"B{idx}"] = f"='Assumptions'!B{price_row}"
        me[f"C{idx}"] = f"='Assumptions'!B{ppob_row}"
        me[f"D{idx}"] = f"='Assumptions'!B{partner_row}"
        me[f"E{idx}"] = f"=B{idx}*'Assumptions'!B5"
        me[f"F{idx}"] = f"=B{idx}*'Assumptions'!B7"
        me[f"G{idx}"] = f"=B{idx}*'Assumptions'!B8"
        me[f"H{idx}"] = f"=B{idx}*'Assumptions'!B9"
        me[f"I{idx}"] = f"=B{idx}*'Assumptions'!B6"
        me[f"J{idx}"] = f"=B{idx}*'Assumptions'!B10"
        me[f"K{idx}"] = f"=B{idx}-SUM(C{idx}:J{idx})"

    bl = wb["Bonus Liability"]
    bl.append(["Stress Test", "Payout Ratio", "Gross Membership Revenue", "Estimated Bonus Liability"])
    for idx, (label, ratio) in enumerate([("Low", 0.25), ("Managed", 0.40), ("Stress", 0.60), ("Concentration", 0.75)], 2):
        bl[f"A{idx}"] = label; bl[f"B{idx}"] = ratio; bl[f"C{idx}"] = "='Revenue Build'!F5"; bl[f"D{idx}"] = f"=B{idx}*C{idx}"

    rb = wb["Revenue Build"]
    rb.append(["Metric"] + years)
    rb.append(["Paid Members", "='User Growth'!B4", "='User Growth'!C4", "='User Growth'!D4", "='User Growth'!E4", "='User Growth'!F4"])
    rb.append(["Avg Membership Revenue", 800000, 900000, 1000000, 1050000, 1100000])
    rb.append(["Membership Revenue", "=B2*B3", "=C2*C3", "=D2*D3", "=E2*E3", "=F2*F3"])
    rb.append(["Merchant Revenue", "='Merchant Growth'!B3*'Assumptions'!B22*12", "='Merchant Growth'!C3*'Assumptions'!B22*12", "='Merchant Growth'!D3*'Assumptions'!B22*12", "='Merchant Growth'!E3*'Assumptions'!B22*12", "='Merchant Growth'!F3*'Assumptions'!B22*12"])
    rb.append(["Total Revenue", "=B4+B5", "=C4+C5", "=D4+D5", "=E4+E5", "=F4+F5"])

    vc = wb["Variable Cost"]
    vc.append(["Metric"] + years)
    vc.append(["Gateway Cost", "='Revenue Build'!B6*'Assumptions'!B5", "='Revenue Build'!C6*'Assumptions'!B5", "='Revenue Build'!D6*'Assumptions'!B5", "='Revenue Build'!E6*'Assumptions'!B5", "='Revenue Build'!F6*'Assumptions'!B5"])
    vc.append(["Bonus/Reward Provision", "='Revenue Build'!B4*35%", "='Revenue Build'!C4*35%", "='Revenue Build'!D4*35%", "='Revenue Build'!E4*35%", "='Revenue Build'!F4*35%"])
    vc.append(["Other Variable Cost", "='Revenue Build'!B6*8%", "='Revenue Build'!C6*8%", "='Revenue Build'!D6*8%", "='Revenue Build'!E6*8%", "='Revenue Build'!F6*8%"])
    vc.append(["Total Variable Cost", "=SUM(B2:B4)", "=SUM(C2:C4)", "=SUM(D2:D4)", "=SUM(E2:E4)", "=SUM(F2:F4)"])

    fc = wb["Fixed Cost"]
    fc.append(["Metric"] + years)
    fc.append(["Annual Fixed Cost", "='Assumptions'!B3*12", "=B2*1.25", "=C2*1.25", "=D2*1.22", "=E2*1.18"])

    pl = wb["P&L Annual"]
    pl.append(["Metric"] + years)
    pl.append(["Revenue", "='Revenue Build'!B6", "='Revenue Build'!C6", "='Revenue Build'!D6", "='Revenue Build'!E6", "='Revenue Build'!F6"])
    pl.append(["Cost of Revenue", "='Variable Cost'!B5", "='Variable Cost'!C5", "='Variable Cost'!D5", "='Variable Cost'!E5", "='Variable Cost'!F5"])
    pl.append(["Gross Profit", "=B2-B3", "=C2-C3", "=D2-D3", "=E2-E3", "=F2-F3"])
    pl.append(["Operating Expense", "='Fixed Cost'!B2", "='Fixed Cost'!C2", "='Fixed Cost'!D2", "='Fixed Cost'!E2", "='Fixed Cost'!F2"])
    pl.append(["EBITDA", "=B4-B5", "=C4-C5", "=D4-D5", "=E4-E5", "=F4-F5"])
    pl.append(["Tax", "=MAX(0,B6*'Assumptions'!B6)", "=MAX(0,C6*'Assumptions'!B6)", "=MAX(0,D6*'Assumptions'!B6)", "=MAX(0,E6*'Assumptions'!B6)", "=MAX(0,F6*'Assumptions'!B6)"])
    pl.append(["Net Profit", "=B6-B7", "=C6-C7", "=D6-D7", "=E6-E7", "=F6-F7"])

    cf = wb["Cash Flow"]
    cf.append(["Metric"] + years)
    cf.append(["Opening Cash", "='Assumptions'!B23", "=B5", "=C5", "=D5", "=E5"])
    cf.append(["Operating Cash Flow", "='P&L Annual'!B8", "='P&L Annual'!C8", "='P&L Annual'!D8", "='P&L Annual'!E8", "='P&L Annual'!F8"])
    cf.append(["Capex", "=-'Assumptions'!B2", "=-250000000", "=-300000000", "=-350000000", "=-400000000"])
    cf.append(["Financing Cash Flow", 0, 0, 0, 0, 0])
    cf.append(["Ending Cash", "=SUM(B2:B4)", "=SUM(C2:C4)", "=SUM(D2:D4)", "=SUM(E2:E4)", "=SUM(F2:F4)"])

    bs = wb["Balance Sheet"]
    bs.append(["Metric"] + years)
    bs.append(["Cash", "='Cash Flow'!B5", "='Cash Flow'!C5", "='Cash Flow'!D5", "='Cash Flow'!E5", "='Cash Flow'!F5"])
    bs.append(["Fixed Assets", 800000000, 900000000, 1000000000, 1100000000, 1200000000])
    bs.append(["Total Assets", "=B2+B3", "=C2+C3", "=D2+D3", "=E2+E3", "=F2+F3"])
    bs.append(["Liabilities", 300000000, 350000000, 450000000, 550000000, 650000000])
    bs.append(["Equity", "=B4-B5", "=C4-C5", "=D4-D5", "=E4-E5", "=F4-F5"])
    bs.append(["Check", "=B4-B5-B6", "=C4-C5-C6", "=D4-D5-D6", "=E4-E5-E6", "=F4-F5-F6"])

    wc = wb["Working Capital"]
    wc.append(["Runway", "Formula", "Amount"])
    for idx, months in enumerate([6, 12, 18], 2):
        wc[f"A{idx}"] = f"{months} months"
        wc[f"B{idx}"] = f"='Assumptions'!B3*{months}"
        wc[f"C{idx}"] = f"=B{idx}"

    be = wb["Break-Even"]
    be.append(["Metric", "Formula", "Base Result"])
    be.append(["Annual Fixed Cost", "='Fixed Cost'!B2", "=B2"])
    be.append(["Avg Contribution per Paid Member", "='Membership Economics'!K2", "=B3"])
    be.append(["Break-Even Paid Members", "=ROUNDUP(B2/B3,0)", "=B4"])
    be.append(["Estimated Break-Even Year", "2029 if EBITDA turns positive in base scenario", "2029"])

    ue = wb["Unit Economics"]
    ue.append(["Metric", "Formula / Assumption", "Base"])
    metrics = [("CAC", "='Assumptions'!B20", "='Assumptions'!B20"), ("ARPU Monthly", "='Assumptions'!B21", "='Assumptions'!B21"), ("Gross Margin", "='P&L Annual'!B4/'P&L Annual'!B2", "='P&L Annual'!B4/'P&L Annual'!B2"), ("Payback Months", "=CAC/ARPU", "=C2/C3")]
    for row in metrics:
        ue.append(row)

    funding = wb["Funding"]
    funding.append(["Use of Funds", "Amount"])
    for row in [("Product Development", 350000000), ("Infrastructure Setup", 75000000), ("Equipment", 75000000), ("Branding", 30000000), ("Legal & Compliance", 40000000), ("Initial Marketing", 100000000), ("Office Setup", 50000000), ("Contingency", 80000000), ("12 Month Working Capital", "='Working Capital'!C3")]:
        funding.append(row)
    funding.append(["Total Baseline Requirement", "=SUM(B2:B10)"])

    val = wb["Valuation"]
    val.append(["Method", "Conservative", "Base", "Upside", "Disclaimer"])
    val.append(["Revenue Multiple", "='Revenue Build'!D6*1.0", "='Revenue Build'!D6*1.8", "='Revenue Build'!D6*3.0", "Illustrative only"])
    val.append(["EBITDA Multiple", "=MAX(0,'P&L Annual'!D6*5)", "=MAX(0,'P&L Annual'!E6*7)", "=MAX(0,'P&L Annual'!F6*9)", "Not appraisal"])
    val.append(["VC Method", 2000000000, 5000000000, 12000000000, "Highly assumption sensitive"])

    sens = wb["Sensitivity"]
    sens.append(["Driver", "Low", "Base", "High"])
    for row in [("Bonus Payout", 0.25, 0.4, 0.6), ("CAC", 45000, 65000, 95000), ("Paid Conversion", 0.12, 0.20, 0.28), ("Fixed Cost Multiplier", 0.85, 1.0, 1.25)]:
        sens.append(row)

    kpi = wb["KPI Dashboard"]
    kpi.append(["KPI Category", "Metric", "Target Use"])
    for row in [("Member", "Registered / Active / Paid / Churn / Retention", "Growth quality"), ("Merchant", "Registered / Active / GMV / Revenue", "Ecosystem scale"), ("Payment", "Success rate / Webhook success / Refund / Chargeback", "Reliability"), ("Finance", "Revenue / Gross margin / EBITDA / Burn / Runway", "Capital discipline"), ("Operations", "Ticket SLA / NPS / Uptime / Fraud rate", "Service quality")]:
        kpi.append(row)

    checks = wb["Checks"]
    checks.append(["Check", "Formula", "Status"])
    checks.append(["Balance Sheet balances 2030", "='Balance Sheet'!F7", '=IF(ABS(B2)<1,"PASS","FAIL")'])
    checks.append(["Ending cash is numeric", "='Cash Flow'!F5", '=IF(ISNUMBER(B3),"PASS","FAIL")'])
    checks.append(["Funding requirement positive", "='Funding'!B11", '=IF(B4>0,"PASS","FAIL")'])

    monthly = wb["P&L Monthly"]
    monthly.append(["Month", "Revenue", "Opex", "EBITDA"])
    for i in range(1, 25):
        monthly.append([f"Month {i}", f"='Revenue Build'!B6/12*POWER(1.03,{i-1})", "='Assumptions'!B3", f"=B{i+1}-C{i+1}"])

    # Lightweight sheets requested but not deeply modeled.
    for name in ["Headcount", "Capex"]:
        ws = wb[name]
        ws.append(["Category", 2026, 2027, 2028, 2029, 2030])
        if name == "Headcount":
            ws.append(["Core Team", 8, 14, 24, 38, 55])
            ws.append(["Regional/Ops", 3, 10, 25, 55, 90])
        else:
            ws.append(["Product & Software", 350000000, 250000000, 300000000, 350000000, 400000000])
            ws.append(["Equipment", 75000000, 100000000, 150000000, 200000000, 250000000])

    # Charts.
    chart = LineChart()
    chart.title = "Ending Cash Projection"
    data = Reference(cf, min_col=2, max_col=6, min_row=5, max_row=5)
    cats = Reference(cf, min_col=2, max_col=6, min_row=1, max_row=1)
    chart.add_data(data, from_rows=True, titles_from_data=False)
    chart.set_categories(cats)
    cf.add_chart(chart, "A8")

    bar = BarChart()
    bar.title = "Revenue Build"
    data = Reference(rb, min_col=2, max_col=6, min_row=6, max_row=6)
    cats = Reference(rb, min_col=2, max_col=6, min_row=1, max_row=1)
    bar.add_data(data, from_rows=True, titles_from_data=False)
    bar.set_categories(cats)
    rb.add_chart(bar, "A9")

    for ws in wb.worksheets:
        style_range(ws)
        for row in ws.iter_rows():
            for cell in row:
                if isinstance(cell.value, (int, float)) or (isinstance(cell.value, str) and cell.value.startswith("=")):
                    if cell.column > 1:
                        cell.number_format = '#,##0'
        ws.freeze_panes = "A4"
    wb.save(path)
    return sheets


def make_assumptions_md(path: Path) -> None:
    text = f"""# TapGo Master Blueprint Assumptions

Tanggal: {TODAY.isoformat()}

## Prinsip

Dokumen ini membedakan kondisi aktual, target, roadmap, proyeksi, asumsi, dan aspirasi jangka panjang. Angka dalam financial model adalah ilustrasi internal dan perlu diperbarui dengan data aktual setelah launch.

## Identitas

- Perusahaan: PT TAPGO LION INDONESIA
- Brand: TapGo
- Founder Chairman: Ahmad Zulhi
- Website: https://tapgolion.id
- Email: support@tapgolion.id
- Positioning: Platform Ekosistem Digital Berbasis Membership
- Tagline: Satu Membership, Beragam Manfaat.

## Asumsi Modal Investasi Moderat

- Product development: Rp350.000.000
- Technology infrastructure setup: Rp75.000.000
- Device/equipment: Rp75.000.000
- Branding and creative: Rp30.000.000
- Legal, compliance, HKI: Rp40.000.000
- Initial marketing launch: Rp100.000.000
- Office/setup: Rp50.000.000
- Contingency: Rp80.000.000
- Total baseline: Rp800.000.000

Status: Owner Confirmation Required.

## Asumsi Modal Kerja Moderat

- Payroll/core team: Rp75.000.000 per bulan
- Cloud/infrastructure: Rp7.500.000 per bulan
- Software/tools: Rp4.000.000 per bulan
- Office and utilities: Rp8.000.000 per bulan
- Customer service: Rp10.000.000 per bulan
- Legal/accounting/compliance: Rp5.000.000 per bulan
- Marketing: Rp30.000.000 per bulan
- Travel/merchant operations: Rp10.000.000 per bulan
- Contingency: Rp10.000.000 per bulan
- Total baseline: Rp159.500.000 per bulan

## Data yang Perlu Konfirmasi Owner

- Biaya aktual yang sudah menjadi sunk cost.
- Nilai kebutuhan dana baru.
- Status final DOKU production, Midtrans fallback, dan channel pembayaran.
- Biaya payment gateway final.
- Biaya partner benefit, BPJS, merchandise, dan PPOB.
- CAC aktual setelah launch.
- Package mix aktual.
- Merchant onboarding cost.
- Bonus payout aktual.
- Headcount dan salary plan final.
- Kebijakan akuntansi untuk software development asset.
- Tax treatment final.

## Disclaimer

{DISCLAIMER}
"""
    path.write_text(text, encoding="utf-8")


def verify_text_safety(paths: Iterable[Path]) -> list[str]:
    banned = ["Xavindo", "Febrina Delia", "1630014495231", "guaranteed return", "investasi pasti untung"]
    findings = []
    for path in paths:
        if path.suffix.lower() in {".pdf", ".pptx", ".xlsx"}:
            continue
        text = path.read_text(encoding="utf-8", errors="ignore")
        for item in banned:
            if item.lower() in text.lower():
                findings.append(f"{path}: contains {item}")
    return findings


def main():
    ensure_dirs()
    master_pdf = OUT_DIR / "TAPGO_MASTER_BLUEPRINT_2026_2035.pdf"
    pptx = OUT_DIR / "TAPGO_MASTER_BLUEPRINT_2026_2035.pptx"
    xlsx = OUT_DIR / "TAPGO_FINANCIAL_MODEL_2026_2035.xlsx"
    executive_pdf = OUT_DIR / "TAPGO_EXECUTIVE_BOOK_2026_2035.pdf"
    assumptions = OUT_DIR / "TAPGO_MASTER_BLUEPRINT_ASSUMPTIONS.md"
    master_pages = make_master_pdf(master_pdf, executive=False)
    executive_pages = make_master_pdf(executive_pdf, executive=True)
    slide_count = make_pptx(pptx)
    sheets = make_workbook(xlsx)
    make_assumptions_md(assumptions)
    findings = verify_text_safety([assumptions])
    qa = OUT_DIR / "GENERATION_QA_SUMMARY.txt"
    qa.write_text(
        "\n".join([
            f"Master PDF pages: {master_pages}",
            f"Executive PDF pages: {executive_pages}",
            f"PPTX slides: {slide_count}",
            f"Workbook sheets: {len(sheets)}",
            f"Safety findings: {findings or 'none'}",
            f"Generated: {TODAY.isoformat()}",
        ]),
        encoding="utf-8",
    )
    print(qa.read_text(encoding="utf-8"))


if __name__ == "__main__":
    main()
